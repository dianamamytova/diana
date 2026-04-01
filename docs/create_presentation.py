from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
BG_DARK = RGBColor(0x1B, 0x1B, 0x2F)
BG_CARD = RGBColor(0x27, 0x27, 0x4A)
ACCENT = RGBColor(0x6C, 0x63, 0xFF)
ACCENT2 = RGBColor(0x00, 0xD2, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
GREEN = RGBColor(0x00, 0xE6, 0x76)


def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_accent_bar(slide, left, top, width, height, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, bullet_color=ACCENT2, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = spacing
        p.level = 0
    return txBox


def add_slide_number(slide, num, total=20):
    add_text_box(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
                 f"{num} / {total}", font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)


# ===================== SLIDE 1 — Cover =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

# Top accent line
add_accent_bar(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT)

# Center content
add_text_box(slide, Inches(1.5), Inches(1.2), Inches(10.3), Inches(0.6),
             "DIPLOMA PROJECT", font_size=16, color=ACCENT2, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(2.0),
             "Development of a Multi-Branch\nCoffee Shop Management Web Application\nUsing Spring Boot",
             font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Divider
add_accent_bar(slide, Inches(5.5), Inches(4.3), Inches(2.3), Inches(0.04), ACCENT)

add_text_box(slide, Inches(1.5), Inches(4.8), Inches(10.3), Inches(0.5),
             "Student: Diana Mamytova  |  Group: COM22", font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(5.4), Inches(10.3), Inches(0.5),
             "Adviser: Nargiza Zhumalieva", font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 1)

# ===================== SLIDE 2 — Introduction =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
             "Introduction", font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.04), ACCENT2)

items = [
    "The hospitality industry is undergoing rapid digital transformation.",
    "Multi-branch coffee shops require centralized, scalable, and secure management systems.",
    "This project proposes a production-ready web solution for digital reservation, menu management, and analytics."
]
add_bullet_list(slide, Inches(1.0), Inches(1.8), Inches(11), Inches(4.5), items, font_size=22, spacing=Pt(18))
add_slide_number(slide, 2)

# ===================== SLIDE 3 — Problem Statement =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
             "Problem Statement", font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.04), ACCENT2)

add_text_box(slide, Inches(1.0), Inches(1.6), Inches(11), Inches(0.8),
             "Many coffee shops still rely on manual processes for reservations and management.",
             font_size=20, color=LIGHT_GRAY)

problems = [
    "Reservation conflicts and overlapping bookings",
    "Lack of centralized branch management",
    "Limited data analytics",
    "No integrated QR-based ecosystem",
    "Poor customer interaction management"
]

y = 2.6
for prob in problems:
    add_shape_bg(slide, Inches(1.0), Inches(y), Inches(11), Inches(0.55), BG_CARD)
    add_text_box(slide, Inches(1.2), Inches(y + 0.05), Inches(0.4), Inches(0.45),
                 "!", font_size=18, color=ORANGE, bold=True)
    add_text_box(slide, Inches(1.7), Inches(y + 0.05), Inches(10), Inches(0.45),
                 prob, font_size=18, color=WHITE)
    y += 0.75

add_slide_number(slide, 3)

# ===================== SLIDE 4 — Research and Market Analysis =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "Research and Market Analysis", font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.04), ACCENT2)

# Existing systems
add_text_box(slide, Inches(1.0), Inches(1.6), Inches(5), Inches(0.5),
             "Existing Systems Analyzed:", font_size=20, color=ACCENT2, bold=True)

systems = ["OpenTable", "Toast POS", "Square for Restaurants"]
y = 2.3
for s in systems:
    add_shape_bg(slide, Inches(1.0), Inches(y), Inches(4.5), Inches(0.55), BG_CARD)
    add_text_box(slide, Inches(1.4), Inches(y + 0.05), Inches(4), Inches(0.45),
                 s, font_size=18, color=WHITE)
    y += 0.7

# Key observations
add_text_box(slide, Inches(7.0), Inches(1.6), Inches(5), Inches(0.5),
             "Key Observations:", font_size=20, color=ACCENT2, bold=True)

observations = ["High subscription costs", "Limited customization", "SaaS dependency", "Restricted data ownership"]
y = 2.3
for o in observations:
    add_shape_bg(slide, Inches(7.0), Inches(y), Inches(5.5), Inches(0.55), BG_CARD)
    add_text_box(slide, Inches(7.4), Inches(y + 0.05), Inches(5), Inches(0.45),
                 o, font_size=18, color=WHITE)
    y += 0.7

add_slide_number(slide, 4)

# ===================== SLIDE 5 — Identified Research Gap =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "Identified Research Gap", font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.04), ACCENT2)

add_text_box(slide, Inches(1.0), Inches(1.6), Inches(11), Inches(1.0),
             "There is no lightweight, fully customizable, production-ready web application\ndesigned specifically for multi-branch coffee shops that:",
             font_size=20, color=LIGHT_GRAY)

gaps = [
    "Provides centralized control",
    "Supports branch-level autonomy",
    "Implements a customizable anti-overlap reservation algorithm",
    "Integrates QR services into a single ecosystem",
    "Ensures full data ownership"
]

y = 3.0
for g in gaps:
    add_accent_bar(slide, Inches(1.2), Inches(y + 0.15), Inches(0.3), Inches(0.04), GREEN)
    add_text_box(slide, Inches(1.8), Inches(y), Inches(10), Inches(0.5),
                 g, font_size=20, color=WHITE)
    y += 0.65

add_slide_number(slide, 5)

# ===================== SLIDE 6 — Project Objective =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "Project Objective", font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.04), ACCENT2)

add_text_box(slide, Inches(1.0), Inches(1.6), Inches(11), Inches(0.6),
             "To design and develop a full-stack web application that:", font_size=20, color=LIGHT_GRAY)

objectives = [
    "Supports multiple coffee shop branches",
    "Provides secure authentication (JWT + Google OAuth)",
    "Prevents reservation conflicts",
    "Includes analytics and statistics",
    "Integrates QR-based functionality",
    "Uses modern production-ready architecture"
]

y = 2.5
for i, obj in enumerate(objectives):
    add_shape_bg(slide, Inches(1.0), Inches(y), Inches(11), Inches(0.55), BG_CARD)
    add_text_box(slide, Inches(1.3), Inches(y + 0.05), Inches(0.4), Inches(0.45),
                 str(i + 1), font_size=16, color=ACCENT2, bold=True)
    add_text_box(slide, Inches(1.8), Inches(y + 0.05), Inches(10), Inches(0.45),
                 obj, font_size=18, color=WHITE)
    y += 0.7

add_slide_number(slide, 6)

# ===================== SLIDE 7 — Technologies Used =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "Technologies Used", font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.04), ACCENT2)

# Backend card
add_shape_bg(slide, Inches(0.6), Inches(1.5), Inches(3.0), Inches(3.0), BG_CARD)
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(2.6), Inches(0.5),
             "Backend", font_size=20, color=ACCENT2, bold=True)
add_bullet_list(slide, Inches(0.9), Inches(2.2), Inches(2.5), Inches(2.0),
                ["Java Spring Boot", "Spring Security", "(JWT + OAuth2)", "Spring Data JPA", "Hibernate"],
                font_size=15, color=WHITE)

# Frontend card
add_shape_bg(slide, Inches(3.9), Inches(1.5), Inches(2.8), Inches(3.0), BG_CARD)
add_text_box(slide, Inches(4.1), Inches(1.6), Inches(2.4), Inches(0.5),
             "Frontend", font_size=20, color=ACCENT2, bold=True)
add_bullet_list(slide, Inches(4.2), Inches(2.2), Inches(2.3), Inches(2.0),
                ["React", "HTML, CSS", "JavaScript"],
                font_size=15, color=WHITE)

# Database card
add_shape_bg(slide, Inches(7.0), Inches(1.5), Inches(2.8), Inches(3.0), BG_CARD)
add_text_box(slide, Inches(7.2), Inches(1.6), Inches(2.4), Inches(0.5),
             "Database", font_size=20, color=ACCENT2, bold=True)
add_bullet_list(slide, Inches(7.3), Inches(2.2), Inches(2.3), Inches(2.0),
                ["PostgreSQL"],
                font_size=15, color=WHITE)

# Tools card
add_shape_bg(slide, Inches(10.1), Inches(1.5), Inches(2.8), Inches(3.0), BG_CARD)
add_text_box(slide, Inches(10.3), Inches(1.6), Inches(2.4), Inches(0.5),
             "Other Tools", font_size=20, color=ACCENT2, bold=True)
add_bullet_list(slide, Inches(10.4), Inches(2.2), Inches(2.3), Inches(2.0),
                ["Swagger (OpenAPI)", "Docker", "CI/CD Pipeline"],
                font_size=15, color=WHITE)

add_slide_number(slide, 7)

# ===================== SLIDE 8 — System Architecture =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "System Architecture", font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.04), ACCENT2)

add_text_box(slide, Inches(1.0), Inches(1.5), Inches(5), Inches(0.5),
             "Architecture Type: Layered Architecture", font_size=20, color=ACCENT2, bold=True)

# Layers
layers = [
    ("Controller Layer", "REST API", ACCENT),
    ("Service Layer", "Business Logic", RGBColor(0x7C, 0x4D, 0xFF)),
    ("Repository Layer", "Database Access", RGBColor(0x00, 0x96, 0xC7)),
]

y = 2.3
for name, desc, clr in layers:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(y), Inches(5), Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = clr
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{name}  —  {desc}"
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    y += 0.9

# Arrow indicators between layers
for ay in [3.0, 3.9]:
    add_text_box(slide, Inches(3.2), Inches(ay), Inches(0.5), Inches(0.4),
                 "v", font_size=20, color=LIGHT_GRAY, bold=True, alignment=PP_ALIGN.CENTER)

# Features
add_text_box(slide, Inches(7.5), Inches(1.5), Inches(5), Inches(0.5),
             "Key Features:", font_size=20, color=ACCENT2, bold=True)

features = [
    "Stateless authentication (JWT)",
    "Role-based access control",
    "REST communication",
    "Dockerized deployment"
]
y = 2.3
for f in features:
    add_shape_bg(slide, Inches(7.5), Inches(y), Inches(5.0), Inches(0.55), BG_CARD)
    add_text_box(slide, Inches(7.9), Inches(y + 0.05), Inches(4.5), Inches(0.45),
                 f, font_size=17, color=WHITE)
    y += 0.7

add_slide_number(slide, 8)

# ===================== SLIDE 9 — User Roles =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "User Roles", font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.04), ACCENT2)

roles = [
    ("USER", GREEN, ["Register and login", "View menu", "Book tables", "Leave reviews"]),
    ("ADMIN", ORANGE, ["Manage tables", "Approve/cancel reservations", "CRUD menu & categories", "Moderate reviews", "View analytics"]),
    ("SUPERADMIN", RGBColor(0xFF, 0x44, 0x44), ["Manage branches", "Manage administrators", "System-level configuration"]),
]

x = 0.6
for role_name, clr, perms in roles:
    card_w = 3.8
    add_shape_bg(slide, Inches(x), Inches(1.5), Inches(card_w), Inches(5.2), BG_CARD)
    add_accent_bar(slide, Inches(x), Inches(1.5), Inches(card_w), Inches(0.06), clr)
    add_text_box(slide, Inches(x + 0.2), Inches(1.7), Inches(card_w - 0.4), Inches(0.5),
                 role_name, font_size=22, color=clr, bold=True)
    add_bullet_list(slide, Inches(x + 0.3), Inches(2.5), Inches(card_w - 0.5), Inches(4.0),
                    perms, font_size=15, color=WHITE, spacing=Pt(10))
    x += 4.1

add_slide_number(slide, 9)

# ===================== SLIDE 10 — Core Functional Modules =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "Core Functional Modules", font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.04), ACCENT2)

modules = [
    "Authentication Module",
    "Multi-Branch Management Module",
    "Reservation Engine (anti-overlap)",
    "Menu & Category Management",
    "Review and Rating System",
    "Analytics Dashboard",
    "QR Ecosystem Module"
]

# Grid layout: 2 rows
positions = [
    (0.6, 1.8), (3.8, 1.8), (7.0, 1.8), (10.2, 1.8),
    (0.6, 4.2), (3.8, 4.2), (7.0, 4.2)
]
colors = [ACCENT, RGBColor(0x7C, 0x4D, 0xFF), GREEN, ORANGE, ACCENT2, RGBColor(0xFF, 0x44, 0x88), RGBColor(0xFF, 0xD7, 0x00)]

for i, mod in enumerate(modules):
    px, py = positions[i]
    add_shape_bg(slide, Inches(px), Inches(py), Inches(2.9), Inches(2.0), BG_CARD)
    add_accent_bar(slide, Inches(px), Inches(py), Inches(2.9), Inches(0.06), colors[i])
    add_text_box(slide, Inches(px + 0.2), Inches(py + 0.5), Inches(2.5), Inches(1.2),
                 mod, font_size=17, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 10)

# ===================== SLIDE 11 — Reservation Engine =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "Reservation Engine", font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.04), ACCENT2)

features = [
    "Time-slot validation",
    "Prevention of overlapping reservations",
    "Reservation editing and cancellation",
    "Administrator approval workflow",
    "Capacity-aware booking logic"
]

y = 1.8
for f in features:
    add_shape_bg(slide, Inches(1.0), Inches(y), Inches(11), Inches(0.6), BG_CARD)
    add_accent_bar(slide, Inches(1.0), Inches(y), Inches(0.08), Inches(0.6), GREEN)
    add_text_box(slide, Inches(1.4), Inches(y + 0.08), Inches(10), Inches(0.45),
                 f, font_size=19, color=WHITE)
    y += 0.8

add_text_box(slide, Inches(1.0), Inches(y + 0.3), Inches(11), Inches(0.5),
             "This ensures operational reliability and fairness.", font_size=18, color=ACCENT2, bold=True)
add_slide_number(slide, 11)

# ===================== SLIDE 12 — QR Ecosystem =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "QR Ecosystem", font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.04), ACCENT2)

qr_items = [
    ("QR code for each table", "Direct menu access"),
    ("QR code for each branch", "Branch information"),
    ("QR code for campaigns", "Promotional content"),
    ("QR-based review", "Customer feedback"),
]

x = 0.6
for title, desc in qr_items:
    add_shape_bg(slide, Inches(x), Inches(1.8), Inches(2.9), Inches(2.8), BG_CARD)
    # QR icon placeholder
    qr_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.8), Inches(2.0), Inches(1.2), Inches(1.2))
    qr_shape.fill.solid()
    qr_shape.fill.fore_color.rgb = ACCENT
    qr_shape.line.fill.background()
    add_text_box(slide, Inches(x + 0.8), Inches(2.2), Inches(1.2), Inches(0.8),
                 "QR", font_size=24, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(x + 0.1), Inches(3.4), Inches(2.7), Inches(0.5),
                 title, font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(x + 0.1), Inches(3.9), Inches(2.7), Inches(0.5),
                 desc, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    x += 3.2

add_text_box(slide, Inches(1.0), Inches(5.2), Inches(11), Inches(0.5),
             "Improves customer experience and reduces operational friction.", font_size=18, color=ACCENT2, bold=True)
add_slide_number(slide, 12)

# ===================== SLIDE 13 — Database Design =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "Database Design (ER Concept)", font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.04), ACCENT2)

entities = ["User", "Branch", "Table", "Reservation", "Category", "MenuItem", "Review"]
entity_colors = [GREEN, ACCENT, ORANGE, RGBColor(0xFF, 0x44, 0x88), ACCENT2, RGBColor(0xFF, 0xD7, 0x00), RGBColor(0x7C, 0x4D, 0xFF)]

# Entity boxes
positions_ent = [
    (0.8, 1.8), (3.4, 1.8), (6.0, 1.8), (8.6, 1.8),
    (2.1, 3.8), (5.0, 3.8), (7.9, 3.8)
]

for i, ent in enumerate(entities):
    px, py = positions_ent[i]
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(px), Inches(py), Inches(2.2), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.color.rgb = entity_colors[i]
    shape.line.width = Pt(2)
    add_text_box(slide, Inches(px + 0.1), Inches(py + 0.3), Inches(2.0), Inches(0.6),
                 ent, font_size=18, color=entity_colors[i], bold=True, alignment=PP_ALIGN.CENTER)

# Relationships text
add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.4),
             "Relationships:", font_size=18, color=ACCENT2, bold=True)
rels = [
    "Branch has many Tables  |  Branch has many MenuItems  |  User creates Reservations  |  User writes Reviews  |  Reservation belongs to Table"
]
add_text_box(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.5),
             rels[0], font_size=14, color=LIGHT_GRAY)
add_slide_number(slide, 13)

# ===================== SLIDE 14 — Analytics and Reporting =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "Analytics and Reporting", font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.04), ACCENT2)

analytics = [
    "Daily reservation statistics",
    "Branch performance comparison",
    "Reservation trends over time",
    "Graph-based data visualization"
]

y = 2.0
for a in analytics:
    add_shape_bg(slide, Inches(1.0), Inches(y), Inches(5.5), Inches(0.65), BG_CARD)
    add_accent_bar(slide, Inches(1.0), Inches(y), Inches(0.08), Inches(0.65), ACCENT2)
    add_text_box(slide, Inches(1.4), Inches(y + 0.1), Inches(5), Inches(0.45),
                 a, font_size=18, color=WHITE)
    y += 0.85

# Chart placeholder
chart_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(1.8), Inches(5.0), Inches(4.5))
chart_shape.fill.solid()
chart_shape.fill.fore_color.rgb = BG_CARD
chart_shape.line.fill.background()

# Simple bar chart representation
bar_heights = [2.5, 3.5, 2.0, 4.0, 3.0]
bar_colors = [ACCENT, ACCENT2, GREEN, ORANGE, ACCENT]
for i, h in enumerate(bar_heights):
    bx = 8.0 + i * 0.8
    by = 5.8 - h * 0.8
    bh = h * 0.8
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(bx), Inches(by), Inches(0.5), Inches(bh))
    bar.fill.solid()
    bar.fill.fore_color.rgb = bar_colors[i]
    bar.line.fill.background()

add_text_box(slide, Inches(1.0), Inches(5.8), Inches(5.5), Inches(0.5),
             "Supports data-driven decision-making.", font_size=18, color=ACCENT2, bold=True)
add_slide_number(slide, 14)

# ===================== SLIDE 15 — MVP =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "Minimum Viable Product (MVP)", font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.04), ACCENT2)

mvp_items = [
    "User registration and login",
    "Branch-based menu display",
    "Online table reservation",
    "Reservation approval by admin",
    "Basic analytics dashboard",
    "QR menu access"
]

y = 1.8
for i, item in enumerate(mvp_items):
    add_shape_bg(slide, Inches(1.0), Inches(y), Inches(11), Inches(0.6), BG_CARD)
    # Checkmark
    add_text_box(slide, Inches(1.3), Inches(y + 0.07), Inches(0.5), Inches(0.45),
                 "v", font_size=16, color=GREEN, bold=True)
    add_text_box(slide, Inches(1.9), Inches(y + 0.07), Inches(9.5), Inches(0.45),
                 item, font_size=18, color=WHITE)
    y += 0.75

add_slide_number(slide, 15)

# ===================== SLIDE 16 — Task Plan =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "Task Plan", font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.04), ACCENT2)

phases = [
    ("Phase 1", "Research and System Design"),
    ("Phase 2", "Database and Backend Development"),
    ("Phase 3", "Frontend Development"),
    ("Phase 4", "Integration and Testing"),
    ("Phase 5", "Dockerization and Deployment"),
    ("Phase 6", "Documentation and Final Report"),
]

phase_colors = [ACCENT, RGBColor(0x7C, 0x4D, 0xFF), ACCENT2, GREEN, ORANGE, RGBColor(0xFF, 0x44, 0x88)]

y = 1.8
for i, (phase, desc) in enumerate(phases):
    # Timeline dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.3), Inches(y + 0.12), Inches(0.3), Inches(0.3))
    dot.fill.solid()
    dot.fill.fore_color.rgb = phase_colors[i]
    dot.line.fill.background()

    # Line connector (except last)
    if i < len(phases) - 1:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.43), Inches(y + 0.42), Inches(0.04), Inches(0.5))
        line.fill.solid()
        line.fill.fore_color.rgb = LIGHT_GRAY
        line.line.fill.background()

    add_text_box(slide, Inches(2.0), Inches(y + 0.05), Inches(2.5), Inches(0.4),
                 phase, font_size=18, color=phase_colors[i], bold=True)
    add_text_box(slide, Inches(4.5), Inches(y + 0.05), Inches(7), Inches(0.4),
                 desc, font_size=18, color=WHITE)
    y += 0.85

add_slide_number(slide, 16)

# ===================== SLIDE 17 — Testing and DevOps =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "Testing and DevOps", font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.04), ACCENT2)

# Three cards
cards = [
    ("Testing", ACCENT, ["Unit testing (JUnit)", "Integration testing"]),
    ("Documentation", ACCENT2, ["Swagger API documentation"]),
    ("Deployment", GREEN, ["Docker containerization", "CI/CD pipeline preparation"]),
]

x = 0.6
for title, clr, items in cards:
    add_shape_bg(slide, Inches(x), Inches(1.8), Inches(3.8), Inches(4.0), BG_CARD)
    add_accent_bar(slide, Inches(x), Inches(1.8), Inches(3.8), Inches(0.06), clr)
    add_text_box(slide, Inches(x + 0.3), Inches(2.1), Inches(3.2), Inches(0.5),
                 title, font_size=22, color=clr, bold=True)
    add_bullet_list(slide, Inches(x + 0.4), Inches(2.9), Inches(3.0), Inches(2.5),
                    items, font_size=16, color=WHITE, spacing=Pt(12))
    x += 4.1

add_slide_number(slide, 17)

# ===================== SLIDE 18 — Value Proposition =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "Value Proposition", font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.04), ACCENT2)

# For Owners
add_shape_bg(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.0), BG_CARD)
add_accent_bar(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(0.06), ORANGE)
add_text_box(slide, Inches(0.9), Inches(1.9), Inches(5.2), Inches(0.5),
             "For Coffee Shop Owners", font_size=22, color=ORANGE, bold=True)
owner_items = [
    "Full control over business data",
    "Reduced booking conflicts",
    "Centralized multi-branch management",
    "Cost-effective alternative to SaaS"
]
add_bullet_list(slide, Inches(1.1), Inches(2.7), Inches(5.0), Inches(3.5),
                owner_items, font_size=17, color=WHITE, spacing=Pt(14))

# For Customers
add_shape_bg(slide, Inches(6.9), Inches(1.6), Inches(5.8), Inches(5.0), BG_CARD)
add_accent_bar(slide, Inches(6.9), Inches(1.6), Inches(5.8), Inches(0.06), GREEN)
add_text_box(slide, Inches(7.2), Inches(1.9), Inches(5.2), Inches(0.5),
             "For Customers", font_size=22, color=GREEN, bold=True)
cust_items = [
    "Convenient online booking",
    "Transparent menu and pricing",
    "Fast QR-based access"
]
add_bullet_list(slide, Inches(7.4), Inches(2.7), Inches(5.0), Inches(3.5),
                cust_items, font_size=17, color=WHITE, spacing=Pt(14))

add_slide_number(slide, 18)

# ===================== SLIDE 19 — Future Work =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "Future Work", font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.04), ACCENT2)

future_items = [
    "Online payment integration",
    "AI-based demand prediction",
    "Customer recommendation system",
    "Mobile application version",
    "Cloud deployment"
]

y = 1.8
for i, item in enumerate(future_items):
    add_shape_bg(slide, Inches(1.0), Inches(y), Inches(11), Inches(0.65), BG_CARD)
    add_accent_bar(slide, Inches(1.0), Inches(y), Inches(0.08), Inches(0.65), ACCENT2)
    add_text_box(slide, Inches(1.4), Inches(y + 0.1), Inches(10), Inches(0.45),
                 item, font_size=19, color=WHITE)
    y += 0.85

add_slide_number(slide, 19)

# ===================== SLIDE 20 — References =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "References (APA Style)", font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.04), ACCENT2)

refs = [
    "Fielding, R. T. (2000). Architectural styles and the design of network-based software architectures. Doctoral dissertation, University of California, Irvine.",
    "Porter, M. E. (1985). Competitive advantage: Creating and sustaining superior performance. Free Press.",
    "OpenTable. (2023). Restaurant reservation management systems.",
    "Square, Inc. (2023). Digital transformation in restaurants.",
    "Toast, Inc. (2023). Restaurant technology solutions."
]

y = 1.6
for ref in refs:
    add_shape_bg(slide, Inches(0.8), Inches(y), Inches(11.5), Inches(0.8), BG_CARD)
    add_text_box(slide, Inches(1.1), Inches(y + 0.1), Inches(11), Inches(0.6),
                 ref, font_size=14, color=LIGHT_GRAY)
    y += 1.0

add_slide_number(slide, 20)

# Save
output_path = os.path.expanduser("/Users/raiymbekdaniiaruulu/IdealProjects/StartTups/claud/diplom/nargiza/Diana_Mamytova_Diploma_Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
