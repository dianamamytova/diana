from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Coffee-themed color scheme
BG = RGBColor(0x1A, 0x0E, 0x08)
BG_CARD = RGBColor(0x2C, 0x18, 0x10)
PRIMARY = RGBColor(0x6F, 0x4E, 0x37)
ACCENT = RGBColor(0xD4, 0xA5, 0x74)
CREAM = RGBColor(0xFF, 0xF1, 0xEC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xBB, 0xAA, 0x99)
GREEN = RGBColor(0x28, 0xA7, 0x45)
RED = RGBColor(0xDC, 0x35, 0x45)
BLUE = RGBColor(0x17, 0xA2, 0xB8)
AMBER = RGBColor(0xFF, 0xC1, 0x07)


def add_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha:
        shape.fill.fore_color.brightness = alpha
    return shape


def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox


def add_bullet_list(slide, left, top, width, height, items, size=16, color=LIGHT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return txBox


def add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


# ===== SLIDE 1: Title =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0), Pt(4), prs.slide_width)
add_text(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.5), "DIPLOMA PROJECT", 14, ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(2.5),
         "Development of a Multi-Branch\nCoffee Shop Management\nWeb Application Using Spring Boot",
         44, WHITE, bold=True, align=PP_ALIGN.CENTER)
add_rect(slide, Inches(5.5), Inches(4.5), Inches(2.3), Pt(3), ACCENT)
add_text(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.5),
         "Student: Diana Mamytova  |  Group: COM22", 18, LIGHT, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
         "Adviser: Nargiza Zhumalieva", 16, LIGHT, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.3), "1 / 20", 11, LIGHT, align=PP_ALIGN.RIGHT)


# ===== SLIDE 2: Introduction =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6), "Introduction", 36, WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(1.5))

add_bullet_list(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(5), [
    "The hospitality industry is undergoing rapid digital transformation.",
    "",
    "Multi-branch coffee shops require centralized, scalable, and secure",
    "management systems for operations.",
    "",
    "This project proposes a production-ready web solution for digital",
    "reservation management, menu administration, and business analytics.",
    "",
    "Built with modern technologies: Java Spring Boot, React, PostgreSQL.",
    "",
    "Deployed at: https://coffee-shop-diana.com"
], 20, LIGHT)
add_text(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.3), "2 / 20", 11, LIGHT, align=PP_ALIGN.RIGHT)


# ===== SLIDE 3: Problem Statement =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6), "Problem Statement", 36, WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2))

add_text(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.5),
         "Many coffee shops still rely on manual processes for reservations and management.", 18, LIGHT)

problems = [
    ("!", "Reservation conflicts and overlapping bookings"),
    ("!", "Lack of centralized branch management"),
    ("!", "Limited data analytics for decision making"),
    ("!", "No integrated QR-based ecosystem"),
    ("!", "Poor customer interaction and feedback management"),
]
y = 2.4
for icon, text in problems:
    add_rect(slide, Inches(0.8), Inches(y), Inches(11), Inches(0.55), BG_CARD)
    add_text(slide, Inches(1.0), Inches(y + 0.05), Inches(0.4), Inches(0.45), icon, 16, AMBER, bold=True)
    add_text(slide, Inches(1.5), Inches(y + 0.05), Inches(10), Inches(0.45), text, 17, WHITE)
    y += 0.7

add_text(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.3), "3 / 20", 11, LIGHT, align=PP_ALIGN.RIGHT)


# ===== SLIDE 4: Research and Market Analysis =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.6), "Research and Market Analysis", 36, WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5))

add_text(slide, Inches(0.8), Inches(1.6), Inches(5), Inches(0.4), "Existing Systems Analyzed:", 20, ACCENT, bold=True)
systems = ["OpenTable", "Toast POS", "Square for Restaurants"]
y = 2.2
for s in systems:
    add_rect(slide, Inches(0.8), Inches(y), Inches(5), Inches(0.55), BG_CARD)
    add_text(slide, Inches(1.2), Inches(y + 0.05), Inches(4), Inches(0.45), s, 17, WHITE)
    y += 0.7

add_text(slide, Inches(7), Inches(1.6), Inches(5), Inches(0.4), "Key Observations:", 20, ACCENT, bold=True)
obs = ["High subscription costs", "Limited customization", "SaaS dependency", "Restricted data ownership"]
y = 2.2
for o in obs:
    add_rect(slide, Inches(7), Inches(y), Inches(5.5), Inches(0.55), BG_CARD)
    add_text(slide, Inches(7.4), Inches(y + 0.05), Inches(5), Inches(0.45), o, 17, WHITE)
    y += 0.7

add_text(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.3), "4 / 20", 11, LIGHT, align=PP_ALIGN.RIGHT)


# ===== SLIDE 5: Identified Research Gap =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.6), "Identified Research Gap", 36, WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2))

add_text(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.8),
         "There is no lightweight, fully customizable, production-ready web application\ndesigned specifically for multi-branch coffee shops that:", 18, LIGHT)

gaps = [
    "Provides centralized control with branch-level autonomy",
    "Implements a customizable anti-overlap reservation algorithm",
    "Integrates QR services into a single ecosystem",
    "Supports role-based access control with multiple administrative levels",
    "Ensures full data ownership and control",
    "Is scalable and deployable using modern DevOps practices (Docker, CI/CD)",
]
y = 2.8
for g in gaps:
    add_text(slide, Inches(1.2), Inches(y), Inches(0.4), Inches(0.35), "—", 18, GREEN, bold=True)
    add_text(slide, Inches(1.7), Inches(y), Inches(10), Inches(0.35), g, 17, WHITE)
    y += 0.5

add_text(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.3), "5 / 20", 11, LIGHT, align=PP_ALIGN.RIGHT)


# ===== SLIDE 6: Project Objective =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.6), "Project Objective", 36, WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(1.8))

add_text(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.5),
         "To design and develop a full-stack web application that:", 20, LIGHT)

objectives = [
    ("1", "Supports multiple coffee shop branches"),
    ("2", "Provides secure authentication (JWT + Google OAuth2)"),
    ("3", "Prevents reservation conflicts with anti-overlap algorithm"),
    ("4", "Includes analytics and statistics dashboard"),
    ("5", "Integrates QR-based functionality for customer access"),
    ("6", "Uses modern production-ready architecture with Docker deployment"),
]
y = 2.4
for num, text in objectives:
    add_rect(slide, Inches(0.8), Inches(y), Inches(11), Inches(0.55), BG_CARD)
    add_text(slide, Inches(1.0), Inches(y + 0.05), Inches(0.5), Inches(0.45), num, 17, ACCENT, bold=True)
    add_text(slide, Inches(1.6), Inches(y + 0.05), Inches(10), Inches(0.45), text, 17, WHITE)
    y += 0.65

add_text(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.3), "6 / 20", 11, LIGHT, align=PP_ALIGN.RIGHT)


# ===== SLIDE 7: Technologies Used =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6), "Technologies Used", 36, WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2))

cols = [
    ("Backend", ACCENT, ["Java 21", "Spring Boot 3.4.1", "Spring Security", "(JWT + OAuth2)", "Spring Data JPA", "Hibernate"]),
    ("Frontend", ACCENT, ["React 18", "JavaScript (ES6+)", "Axios", "Chart.js", "React Router 6"]),
    ("Database", ACCENT, ["PostgreSQL 16"]),
    ("DevOps & Tools", ACCENT, ["Docker", "Nginx", "Let's Encrypt SSL", "Swagger (OpenAPI)", "Cloudflare DNS"]),
]
x = 0.8
for title, color, items in cols:
    w = 2.8
    add_rect(slide, Inches(x), Inches(1.6), Inches(w), Inches(5), BG_CARD)
    add_text(slide, Inches(x + 0.2), Inches(1.7), Inches(w), Inches(0.4), title, 18, ACCENT, bold=True)
    y = 2.3
    for item in items:
        add_text(slide, Inches(x + 0.3), Inches(y), Inches(w - 0.4), Inches(0.35), item, 15, WHITE)
        y += 0.4
    x += 3.05

add_text(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.3), "7 / 20", 11, LIGHT, align=PP_ALIGN.RIGHT)


# ===== SLIDE 8: System Architecture =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.6), "System Architecture", 36, WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2))

add_text(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.4), "Architecture Type: Layered Architecture", 18, ACCENT, bold=True)

# Layers
layers = [
    ("Controller Layer — REST API", RGBColor(0x6F, 0x4E, 0x37)),
    ("Service Layer — Business Logic", RGBColor(0x8B, 0x6B, 0x5A)),
    ("Repository Layer — Database Access", RGBColor(0xA0, 0x78, 0x5A)),
]
y = 2.2
for text, color in layers:
    add_rect(slide, Inches(0.8), Inches(y), Inches(5.5), Inches(0.7), color)
    add_text(slide, Inches(1.2), Inches(y + 0.1), Inches(5), Inches(0.5), text, 17, WHITE, bold=True, align=PP_ALIGN.CENTER)
    if y < 3.5:
        add_text(slide, Inches(3.2), Inches(y + 0.65), Inches(1), Inches(0.3), "V", 14, LIGHT, align=PP_ALIGN.CENTER)
    y += 0.95

# Key features
add_text(slide, Inches(7), Inches(1.5), Inches(5), Inches(0.4), "Key Features:", 18, ACCENT, bold=True)
features = ["Stateless authentication (JWT)", "Role-based access control", "REST communication (JSON)", "Dockerized deployment", "SSL/HTTPS encryption"]
y = 2.2
for f in features:
    add_rect(slide, Inches(7), Inches(y), Inches(5.5), Inches(0.5), BG_CARD)
    add_text(slide, Inches(7.3), Inches(y + 0.05), Inches(5), Inches(0.4), f, 16, WHITE)
    y += 0.6

add_text(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.3), "8 / 20", 11, LIGHT, align=PP_ALIGN.RIGHT)


# ===== SLIDE 9: User Roles =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6), "User Roles", 36, WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(1.5))

roles = [
    ("USER", GREEN, ["Register and login", "View menu and branches", "Book tables", "Leave reviews", "View personal dashboard"]),
    ("ADMIN", AMBER, ["Manage tables and menu", "Approve/cancel reservations", "CRUD menu & categories", "Moderate reviews", "View analytics"]),
    ("SUPERADMIN", RED, ["Manage branches", "Manage administrators", "System-level configuration", "All admin permissions"]),
]
x = 0.8
for title, color, perms in roles:
    add_rect(slide, Inches(x), Inches(1.5), Inches(3.7), Inches(5.2), BG_CARD)
    # Colored top line
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.5), Inches(3.7), Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    add_text(slide, Inches(x + 0.3), Inches(1.7), Inches(3), Inches(0.4), title, 20, color, bold=True)
    y = 2.3
    for p in perms:
        add_text(slide, Inches(x + 0.3), Inches(y), Inches(3.2), Inches(0.35), p, 15, WHITE)
        y += 0.45
    x += 4.0

add_text(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.3), "9 / 20", 11, LIGHT, align=PP_ALIGN.RIGHT)


# ===== SLIDE 10: Core Functional Modules =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.6), "Core Functional Modules", 36, WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5))

modules = [
    "Authentication Module", "Multi-Branch\nManagement Module",
    "Reservation Engine\n(anti-overlap)", "Menu & Category\nManagement",
    "Review and Rating\nSystem", "Analytics Dashboard",
    "QR Ecosystem Module", "File Upload Module"
]
colors = [ACCENT, RGBColor(0x8B, 0x5C, 0xF6), GREEN, AMBER, BLUE, RGBColor(0xEC, 0x48, 0x99), RGBColor(0xF5, 0x9E, 0x0B), PRIMARY]
x, y = 0.6, 1.6
for i, (mod, col) in enumerate(zip(modules, colors)):
    add_rect(slide, Inches(x), Inches(y), Inches(2.8), Inches(1.6), BG_CARD)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(2.8), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = col
    line.line.fill.background()
    add_text(slide, Inches(x + 0.2), Inches(y + 0.4), Inches(2.4), Inches(1), mod, 16, WHITE, bold=True, align=PP_ALIGN.CENTER)
    x += 3.05
    if (i + 1) % 4 == 0:
        x = 0.6
        y += 2.0

add_text(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.3), "10 / 20", 11, LIGHT, align=PP_ALIGN.RIGHT)


# ===== SLIDE 11: Reservation Engine =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6), "Reservation Engine", 36, WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2))

features = [
    "Time-slot validation with date and time range checks",
    "Prevention of overlapping reservations (anti-overlap SQL query)",
    "Guest count vs table capacity validation",
    "Reservation editing and cancellation by users",
    "Administrator confirmation workflow (PENDING → CONFIRMED)",
    "Status management: PENDING | CONFIRMED | CANCELLED | COMPLETED",
]
y = 1.8
for f in features:
    add_rect(slide, Inches(0.8), Inches(y), Inches(11.5), Inches(0.55), BG_CARD)
    add_text(slide, Inches(1.0), Inches(y + 0.05), Inches(0.3), Inches(0.4), "|", 18, GREEN, bold=True)
    add_text(slide, Inches(1.4), Inches(y + 0.05), Inches(10.5), Inches(0.4), f, 17, WHITE)
    y += 0.7

add_text(slide, Inches(0.8), Inches(6.0), Inches(11), Inches(0.4),
         "This ensures operational reliability and fairness.", 18, ACCENT, bold=True)

add_text(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.3), "11 / 20", 11, LIGHT, align=PP_ALIGN.RIGHT)


# ===== SLIDE 12: QR Ecosystem =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6), "QR Ecosystem", 36, WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(1.5))

qr_types = [
    ("QR", "QR code for each table", "Direct menu access"),
    ("QR", "QR code for each branch", "Branch information"),
    ("QR", "QR code for campaigns", "Promotional content"),
    ("QR", "QR-based review", "Customer feedback"),
]
x = 0.5
for icon, title, desc in qr_types:
    add_rect(slide, Inches(x), Inches(1.8), Inches(2.9), Inches(3.5), BG_CARD)
    # QR icon box
    add_rect(slide, Inches(x + 0.7), Inches(2.1), Inches(1.5), Inches(1.5), RGBColor(0x6C, 0x63, 0xFF))
    add_text(slide, Inches(x + 0.7), Inches(2.4), Inches(1.5), Inches(0.8), icon, 28, WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(x + 0.2), Inches(3.9), Inches(2.5), Inches(0.4), title, 15, WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(x + 0.2), Inches(4.4), Inches(2.5), Inches(0.4), desc, 14, LIGHT, align=PP_ALIGN.CENTER)
    x += 3.15

add_text(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(0.4),
         "Improves customer experience and reduces operational friction.", 18, ACCENT, bold=True)

add_text(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.3), "12 / 20", 11, LIGHT, align=PP_ALIGN.RIGHT)


# ===== SLIDE 13: Database Design =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.6), "Database Design (ER Concept)", 36, WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5))

entities = [
    ("User", GREEN), ("Branch", RGBColor(0x8B, 0x5C, 0xF6)),
    ("Table", AMBER), ("Reservation", RGBColor(0xEC, 0x48, 0x99)),
    ("Category", BLUE), ("MenuItem", RGBColor(0xF5, 0x9E, 0x0B)),
    ("Review", RGBColor(0x8B, 0x5C, 0xF6)),
]
x, y = 0.5, 1.6
for i, (name, color) in enumerate(entities):
    add_rect(slide, Inches(x), Inches(y), Inches(2.5), Inches(1.3), BG_CARD)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(2.5), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    add_text(slide, Inches(x), Inches(y + 0.35), Inches(2.5), Inches(0.5), name, 20, color, bold=True, align=PP_ALIGN.CENTER)
    x += 3.0
    if (i + 1) % 4 == 0:
        x = 1.8
        y += 1.8

add_text(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.3), "Relationships:", 18, ACCENT, bold=True)
add_text(slide, Inches(0.8), Inches(5.9), Inches(11), Inches(0.4),
         "Branch has many Tables  |  Branch has many MenuItems  |  User creates Reservations  |  User writes Reviews  |  Reservation belongs to Table",
         14, LIGHT)

add_text(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.3), "13 / 20", 11, LIGHT, align=PP_ALIGN.RIGHT)


# ===== SLIDE 14: Analytics and Reporting =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.6), "Analytics and Reporting", 36, WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2))

features = [
    "Daily reservation statistics per branch",
    "Branch performance comparison (confirmed vs cancelled)",
    "Reservation trends over time (line charts)",
    "Rating distribution analysis (5-star breakdown)",
    "Graph-based data visualization using Chart.js",
]
y = 1.8
for f in features:
    add_rect(slide, Inches(0.8), Inches(y), Inches(6), Inches(0.55), BG_CARD)
    add_text(slide, Inches(1.0), Inches(y + 0.05), Inches(0.3), Inches(0.4), "|", 18, ACCENT, bold=True)
    add_text(slide, Inches(1.4), Inches(y + 0.05), Inches(5), Inches(0.4), f, 17, WHITE)
    y += 0.7

# Chart placeholder
add_rect(slide, Inches(7.5), Inches(1.8), Inches(4.8), Inches(4), BG_CARD)
bars_data = [2.0, 3.2, 1.5, 3.8, 2.8]
bar_colors = [BLUE, ACCENT, GREEN, AMBER, RGBColor(0x8B, 0x5C, 0xF6)]
for i, (h, c) in enumerate(zip(bars_data, bar_colors)):
    add_rect(slide, Inches(8.0 + i * 0.85), Inches(5.3 - h), Inches(0.6), Inches(h), c)

add_text(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(0.4),
         "Supports data-driven decision-making.", 18, ACCENT, bold=True)

add_text(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.3), "14 / 20", 11, LIGHT, align=PP_ALIGN.RIGHT)


# ===== SLIDE 15: Authentication & Security =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.6), "Authentication & Security", 36, WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2))

add_text(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.4), "JWT Token Authentication:", 18, ACCENT, bold=True)
jwt_steps = [
    "1. User submits credentials",
    "2. Server validates & generates JWT (HS256)",
    "3. Token stored in localStorage",
    "4. Axios attaches Bearer token to all requests",
    "5. Server validates token on each request",
]
y = 2.2
for s in jwt_steps:
    add_text(slide, Inches(1.0), Inches(y), Inches(5), Inches(0.3), s, 15, WHITE)
    y += 0.4

add_text(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(0.4), "Google OAuth2 Flow:", 18, ACCENT, bold=True)
oauth_steps = [
    "1. Click 'Sign in with Google'",
    "2. Redirect to Google authorization",
    "3. Google callback with auth code",
    "4. Server exchanges code for user info",
    "5. JWT generated, redirect to app",
]
y = 2.2
for s in oauth_steps:
    add_text(slide, Inches(7.2), Inches(y), Inches(5), Inches(0.3), s, 15, WHITE)
    y += 0.4

add_text(slide, Inches(0.8), Inches(4.8), Inches(11.5), Inches(0.4), "Security Features:", 18, ACCENT, bold=True)
sec = ["BCrypt password hashing", "CORS protection", "CSRF disabled (stateless)", "HTTPS/SSL encryption", "Role-based access control"]
x = 0.8
for s in sec:
    add_rect(slide, Inches(x), Inches(5.3), Inches(2.2), Inches(0.5), BG_CARD)
    add_text(slide, Inches(x + 0.1), Inches(5.33), Inches(2), Inches(0.4), s, 12, WHITE, align=PP_ALIGN.CENTER)
    x += 2.35

add_text(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.3), "15 / 20", 11, LIGHT, align=PP_ALIGN.RIGHT)


# ===== SLIDE 16: API Overview =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.6), "REST API Overview", 36, WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2))

add_text(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.4),
         "46 RESTful API endpoints across 10 controllers:", 18, LIGHT)

api_groups = [
    ("Authentication", "3 endpoints", "Register, Login, Profile"),
    ("Branches", "5 endpoints", "CRUD + listing"),
    ("Tables", "5 endpoints", "CRUD + availability"),
    ("Categories", "6 endpoints", "CRUD + by branch"),
    ("Menu Items", "7 endpoints", "CRUD + filtering"),
    ("Reservations", "7 endpoints", "Create, status, cancel"),
    ("Reviews", "6 endpoints", "Submit, moderate, list"),
    ("Analytics", "2 endpoints", "Branch stats, daily"),
    ("QR Codes", "3 endpoints", "Table, branch, review"),
    ("File Upload", "2 endpoints", "Upload, serve"),
]
y = 2.2
for i, (name, count, desc) in enumerate(api_groups):
    x = 0.8 if i < 5 else 7.0
    yy = y if i < 5 else y - 4.5 * 0.5
    add_rect(slide, Inches(x), Inches(2.2 + (i % 5) * 0.8), Inches(5.5), Inches(0.65), BG_CARD)
    add_text(slide, Inches(x + 0.2), Inches(2.25 + (i % 5) * 0.8), Inches(2), Inches(0.3), name, 15, ACCENT, bold=True)
    add_text(slide, Inches(x + 2.2), Inches(2.25 + (i % 5) * 0.8), Inches(1.2), Inches(0.3), count, 13, GREEN)
    add_text(slide, Inches(x + 3.3), Inches(2.25 + (i % 5) * 0.8), Inches(2), Inches(0.3), desc, 13, LIGHT)

add_text(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.3), "16 / 20", 11, LIGHT, align=PP_ALIGN.RIGHT)


# ===== SLIDE 17: Deployment & DevOps =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.6), "Deployment & DevOps", 36, WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2))

cols = [
    ("Docker", RGBColor(0x00, 0x96, 0xE6), ["3 containers:", "• PostgreSQL 16", "• Spring Boot (JDK 21)", "• Nginx + React build", "", "docker-compose.yml", "orchestration"]),
    ("SSL & DNS", GREEN, ["Let's Encrypt SSL", "Auto-renewal (Certbot)", "", "Cloudflare DNS", "• Domain management", "• DDoS protection", "", "HTTPS everywhere"]),
    ("Infrastructure", AMBER, ["VPS: Ubuntu 24.04", "RAM: 8 GB", "SSD: 145 GB", "", "Domain:", "coffee-shop-diana.com", "", "Nginx reverse proxy"]),
]
x = 0.6
for title, color, items in cols:
    add_rect(slide, Inches(x), Inches(1.6), Inches(3.8), Inches(5), BG_CARD)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.6), Inches(3.8), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    add_text(slide, Inches(x + 0.3), Inches(1.8), Inches(3.2), Inches(0.4), title, 20, color, bold=True)
    y = 2.4
    for item in items:
        add_text(slide, Inches(x + 0.3), Inches(y), Inches(3.2), Inches(0.3), item, 14, WHITE if item else LIGHT)
        y += 0.35
    x += 4.1

add_text(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.3), "17 / 20", 11, LIGHT, align=PP_ALIGN.RIGHT)


# ===== SLIDE 18: Value Proposition =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.6), "Value Proposition", 36, WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2))

# Left column
add_rect(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5), BG_CARD)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.5), Pt(4))
line.fill.solid()
line.fill.fore_color.rgb = AMBER
line.line.fill.background()
add_text(slide, Inches(1.2), Inches(1.7), Inches(5), Inches(0.4), "For Coffee Shop Owners", 20, AMBER, bold=True)
owner_items = ["Full control over business data", "Elimination of reservation conflicts", "Centralized management of multiple branches", "Cost-effective alternative to SaaS platforms", "Data-driven decisions through analytics"]
y = 2.3
for item in owner_items:
    add_text(slide, Inches(1.2), Inches(y), Inches(5), Inches(0.3), item, 16, WHITE)
    y += 0.45

# Right column
add_rect(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(5), BG_CARD)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7), Inches(1.5), Inches(5.5), Pt(4))
line.fill.solid()
line.fill.fore_color.rgb = GREEN
line.line.fill.background()
add_text(slide, Inches(7.4), Inches(1.7), Inches(5), Inches(0.4), "For Customers", 20, GREEN, bold=True)
cust_items = ["Convenient online booking", "Transparent menu with pricing", "Fast QR-based access to services", "Reliable booking confirmation", "Easy review submission"]
y = 2.3
for item in cust_items:
    add_text(slide, Inches(7.4), Inches(y), Inches(5), Inches(0.3), item, 16, WHITE)
    y += 0.45

add_text(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.3), "18 / 20", 11, LIGHT, align=PP_ALIGN.RIGHT)


# ===== SLIDE 19: Future Work =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6), "Future Work", 36, WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(1.5))

future = [
    "Online payment integration (Stripe, PayPal)",
    "AI-based demand prediction for optimal staffing",
    "Customer recommendation system based on order history",
    "Mobile application version (React Native / Flutter)",
    "Cloud deployment with auto-scaling (AWS / GCP)",
    "Email/SMS notifications for reservation confirmations",
]
y = 1.8
for f in future:
    add_rect(slide, Inches(0.8), Inches(y), Inches(11.5), Inches(0.6), BG_CARD)
    add_text(slide, Inches(1.0), Inches(y + 0.08), Inches(0.3), Inches(0.4), "|", 18, ACCENT, bold=True)
    add_text(slide, Inches(1.4), Inches(y + 0.08), Inches(10.5), Inches(0.4), f, 17, WHITE)
    y += 0.75

add_text(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.3), "19 / 20", 11, LIGHT, align=PP_ALIGN.RIGHT)


# ===== SLIDE 20: Thank You =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(1), Inches(2), Inches(11), Inches(1.5),
         "Thank You", 56, WHITE, bold=True, align=PP_ALIGN.CENTER)
add_rect(slide, Inches(5.5), Inches(3.5), Inches(2.3), Pt(3), ACCENT)
add_text(slide, Inches(1), Inches(4), Inches(11), Inches(0.5),
         "Questions & Discussion", 24, ACCENT, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.5),
         "Live Demo: https://coffee-shop-diana.com", 18, LIGHT, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
         "Diana Mamytova  |  COM22  |  Ala-Too International University", 16, LIGHT, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.3), "20 / 20", 11, LIGHT, align=PP_ALIGN.RIGHT)


# Save
output_path = os.path.join(os.path.dirname(__file__), "CoffeeHub_Diploma_Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
